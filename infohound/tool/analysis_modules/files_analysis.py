import exiftool
import requests
import json
import trio
import httpx
import time
import os
import textract
from random import randint
from django.db import IntegrityError
from infohound.models import Domain,Files, Emails
from infohound.tool.retriever_modules import emails as emails_utils
import infohound.tool.infohound_utils as infohound_utils

import sys
import rarfile
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PyPDF2 import PdfReader
from svglib.svglib import svg2rlg
from svgwrite import Drawing


download_direcotry = "infohound/tool/downloaded_files/"
rate_limiter = trio.CapacityLimiter(2)


async def download_file(url, filepath):
    max_retries = 5
    file_url = url
    for i in range(max_retries + 1):
        try:
            async with httpx.AsyncClient(timeout=httpx.Timeout(60.0)) as client:
                user_agent = infohound_utils.getUserAgents()
                response = await client.get(file_url, headers=user_agent[randint(0, len(user_agent)-1)])
                with open(filepath, 'wb') as f:
                    f.write(response.content)
                    print(f"File {url} downloaded!")
                    return 
        except Exception as e:
            print(e)
            print(type(e))
            print(f"Download failed for {url}. Retrying... ({i}/{max_retries})")
            if i%1 == 0:
                file_url = url.replace("https://", "http://")
            else:
                file_url = url              
            time.sleep(5)

async def download_all_files(domain_id):
    queryset = Files.objects.filter(metadata__isnull=True,domain_id=domain_id)
    #async with trio.open_nursery() as nursery:
    for entry in queryset.iterator():
        url = entry.url_download
        filename = entry.filename
        filepath = download_direcotry+filename
        if not os.path.isfile(os.path.join(download_direcotry, filename)):
            #async with rate_limiter:
            #task = nursery.start_soon(download_file, url, filepath)
            await download_file(url,filepath)


def downloadSingleFile(url, filename):
    print("Trying to download: " + filename)
    retry = 0
    downloaded = False
    filepath = download_direcotry+filename
    while retry < 5 and not downloaded:
        try:
            res = requests.get(url)
            open(filepath, 'wb').write(res.content)
            downloaded = True
        except Exception as e:
            retry = retry + 1
        if not downloaded:
            print(retry)
    return downloaded


def getMetadata(domain_id):
    queryset = Files.objects.filter(metadata__isnull=True,domain_id=domain_id)
    for entry in queryset.iterator():
        url = entry.url
        filename = entry.filename
        if os.path.isfile(os.path.join(download_direcotry, filename)):
            extracted = False
            retry = 0
            while not extracted and retry < 5:
                with exiftool.ExifToolHelper() as et:
                    try :
                        filepath = download_direcotry+filename
                        metadata = et.get_metadata([filepath])[0]
                        entry.metadata = metadata
                        extracted = True
                        entry.save()
                        print("metadata extracted")
                    except Exception as e:
                        print(e)
                        print(type(e))
                        retry = retry + 1
        
        


def getEmailsFromMetadata(domain_id):
    queryset = Files.objects.filter(metadata__isnull=False,domain_id=domain_id)
    for entry in queryset.iterator():
        filename = entry.filename
        metadata = str(entry.metadata)
        emails = emails_utils.getEmailsFromText(metadata)
        for e in emails:
            b,em = emails_utils.isValidEmail(e)
            if b:
                domain = Domain.objects.get(id=domain_id).domain
                if domain in em:
                    try:
                        Emails.objects.get_or_create(email=em,source="Files",domain_id=domain_id)
                    except IntegrityError as e:
                        pass



def getEmailsFromFilesContent(domain_id):
    excluded = ["rar","zip"]
    queryset = Files.objects.filter(domain_id=domain_id)
    for entry in queryset.iterator():
        url = entry.url
        filename = entry.filename
        if os.path.isfile(os.path.join(download_direcotry, filename)):
            ext = filename.split(".")[-1:][0]
            if ext not in excluded:
                #text = textract.process(os.path.join(download_direcotry, filename))
                text = extract_text(os.path.join(download_direcotry, filename))
                emails = emails_utils.getEmailsFromText(text)
                for e in emails:
                    b,em = emails_utils.isValidEmail(e)
                    if b:
                        domain = Domain.objects.get(id=domain_id).domain
                        if domain in em:
                            print("Found another email: " + em)
                            try:
                                Emails.objects.get_or_create(email=em,source="Files",domain_id=domain_id)
                            except IntegrityError as e:
                                pass
    







# Currently using textract
def extract_text(file_path):
    text = ""
    file_extension = file_path.split(".")[-1:][0]
    try:
        if file_extension in ["doc", "docx"]:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_extension in ["ppt", "pptx", "pps", "ppsx"]:
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"

        elif file_extension in ["xls", "xlsx"]:
            wb = load_workbook(file_path)
            for sheet in wb:
                for row in sheet.iter_rows():
                    for cell in row:
                        text += str(cell.value) + " "

        elif file_extension == "pdf":
            pdf = PdfReader(file_path)
            for page_num in range(len(pdf.pages)):
                text += pdf.pages[page_num].extract_text()

        elif file_extension == "svg":
            svg = svg2rlg(file_path)
            text = str(svg)

        elif file_extension == "indd":
            print("InDesign file format (.indd) is not supported in this script.")
        
        elif file_extension == "rdp" or file_extension == "ica":
            with open(file_path, 'r') as f:
                text = f.read()

        elif file_extension == "rar":
            with rarfile.RarFile(file_path, 'r') as rar_ref:
                for file in rar_ref.namelist():
                    if not os.path.isdir(file):
                        with rar_ref.open(file) as f:
                            text += f.read().decode('utf-8', errors='ignore')

        else:
            print("Unsupported file format.")
    except Exception as e:
        print("Error occured")

    return text
